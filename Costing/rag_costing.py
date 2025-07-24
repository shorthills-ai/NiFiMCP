import os
import logging
from io import BytesIO
from typing import List, Dict, Any
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, Form, Body
from fastapi.responses import JSONResponse
from fastapi import FastAPI, HTTPException, Request
from pydantic import BaseModel
import numpy as np
import httpx
from openai import AsyncAzureOpenAI
import weaviate
from weaviate.auth import AuthApiKey
from weaviate.classes.query import MetadataQuery
from fastapi import Request
import json
from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Dict, Any
import weaviate
from weaviate.auth import AuthApiKey
from weaviate.classes.query import MetadataQuery
import uvicorn
import traceback
import logging
import json
import os
import time
import numpy as np
from fastapi import Header
from openai import AsyncAzureOpenAI
from httpx_aiohttp import AiohttpTransport
from aiohttp import ClientSession
import openai
import grpc
from dotenv import load_dotenv
import tiktoken
from pydantic import BaseModel, Field
from typing import Optional
from pydantic import BaseModel, Field
from typing import List, Dict, Optional
load_dotenv()
# Setup basic logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
# === Load Env Variables ===
#ACCESS_TOKEN = os.getenv("ACCESS_TOKEN")
DRIVE_ID = os.getenv("DRIVE_ID")
WEAVIATE_URL = os.getenv("WEAVIATE_URL")
WEAVIATE_API_KEY = os.getenv("WEAVIATE_API_KEY")
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "text-embedding-ada-002")
GRAPH_BASE = "https://graph.microsoft.com/v1.0"
#HEADERS = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Accept": "application/json"}
import re
# === Logging Setup ===
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("combined_service")
CHAT_MODEL = "gpt-4o-mini"

# === FastAPI App ===
app = FastAPI(title="Combined RAG Service")

# ========== Service 1: Text Extraction ==========

def extract_text_from_docx(file_bytes: BytesIO) -> List[str]:
    doc = Document(file_bytes)
    chunks = []
    current_chunk = []
    current_size = 0
    max_chunk_size = 1028
    min_chunk_size = 200

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if current_size + len(text) > max_chunk_size and current_chunk:
            chunk_text = " ".join(current_chunk)
            if len(chunk_text) >= min_chunk_size:
                chunks.append(chunk_text)
            current_chunk = []
            current_size = 0

        current_chunk.append(text)
        current_size += len(text)

    if current_chunk:
        chunk_text = " ".join(current_chunk)
        if len(chunk_text) >= min_chunk_size:
            chunks.append(chunk_text)

    return chunks

def extract_text_from_pptx(file_bytes: BytesIO) -> List[str]:
    prs = Presentation(file_bytes)
    chunks = []
    current_chunk = []
    current_size = 0
    max_chunk_size = 1028
    min_chunk_size = 100

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_texts.extend(line.strip() for line in text.split("\n") if line.strip())
        slide_content = " ".join(slide_texts)
        slide_size = len(slide_content)

        if current_size + slide_size > max_chunk_size and current_chunk:
            chunk_text = " ".join(current_chunk)
            if len(chunk_text) >= min_chunk_size:
                chunks.append(chunk_text)
            current_chunk = []
            current_size = 0

        current_chunk.append(slide_content)
        current_size += slide_size

        if current_size >= max_chunk_size:
            chunk_text = " ".join(current_chunk)
            if len(chunk_text) >= min_chunk_size:
                chunks.append(chunk_text)
            current_chunk = []
            current_size = 0

    if current_chunk:
        chunk_text = " ".join(current_chunk)
        if len(chunk_text) >= min_chunk_size:
            chunks.append(chunk_text)

    return chunks



@app.post("/extract-text")
async def extract_text(request: Request, file: UploadFile = File(...)):
    try:
        # === Log Headers ===
        headers = dict(request.headers)
        logger.info(f"Received request headers: {headers}")

        # Extract web_url from headers (case-insensitive match)
        web_url = headers.get("web_url") or headers.get("Web-Url") or None
        
        logger.info(f"Extracted web_url from headers: {web_url}")

        # === Read file ===
        file_bytes = BytesIO(await file.read())
        filename = file.filename
        logger.info(f"Received file: {filename}")

        filename_lower = filename.lower()

        if filename_lower.endswith(".docx"):
            chunks = extract_text_from_docx(file_bytes)
            file_type = "docx"
        elif filename_lower.endswith(".pptx"):
            chunks = extract_text_from_pptx(file_bytes)
            file_type = "pptx"
        else:
            logger.error("Unsupported file type")
            return JSONResponse(status_code=400, content={"error": "Unsupported file type"})

        # === Handle fallback for web_url ===
        if not web_url:
            logger.warning("No webUrl provided in headers; defaulting to filename")
            web_url = filename
        else:
            logger.info(f"Using provided webUrl: {web_url}")

        # === Process Chunks ===
        chunk_output = []
        for idx, chunk in enumerate(chunks):
            cleaned_chunk = " ".join(chunk.split())
            cleaned_chunk = "".join(char for char in cleaned_chunk if ord(char) < 128)

            chunk_metadata = {
                "content": cleaned_chunk,
                "metadata": {
                    "source_file": filename,
                    "file_type": file_type,
                    "chunk_index": idx,
                    "filepath": web_url
                }
            }
            logger.debug(f"Processed chunk {idx}: {chunk_metadata}")
            chunk_output.append(chunk_metadata)

        response_data = {"chunks": chunk_output}
        response_json = json.dumps(response_data, ensure_ascii=True, separators=(',', ':'))

        logger.info(f"Final JSON response preview: {response_json[:200]}")

        return JSONResponse(
            content=json.loads(response_json),
            media_type="application/json"
        )

    except Exception as e:
        logger.exception("Text extraction failed")
        return JSONResponse(status_code=500, content={"error": str(e)})

# ========== Service 2: Embedding ==========

class ChunkInput(BaseModel):
    content: str
    metadata: Dict[str, Any]

class EmbeddingRequest(BaseModel):
    chunks: List[ChunkInput]

async def generate_embedding(text: str):
    try:
        async with httpx.AsyncClient() as client:
            openai_client = AsyncAzureOpenAI(
                azure_endpoint=AZURE_OPENAI_ENDPOINT,
                api_key=AZURE_OPENAI_API_KEY,
                api_version=AZURE_OPENAI_API_VERSION,
                http_client=client
            )
            logger.info(f"Generating embedding for: {text[:80]}...")
            response = await openai_client.embeddings.create(
                model=EMBEDDING_MODEL,
                input=[text],
                encoding_format="float"
            )
            return np.array(response.data[0].embedding)
    except Exception as e:
        logger.error(f"Embedding error: {e}")
        return None


def clean_json_string(json_str: str) -> str:
    """Clean and normalize JSON string."""
    # Remove any leading/trailing whitespace
    json_str = json_str.strip()
    
    # Find the first '{' and last '}'
    start = json_str.find('{')
    end = json_str.rfind('}')
    
    if start != -1 and end != -1:
        # Extract just the JSON object
        json_str = json_str[start:end + 1]
    
    # Remove whitespace between key and colon
    json_str = re.sub(r'"\s+:', '":', json_str)
    
    # Remove any BOM characters
    json_str = json_str.strip('\ufeff')
    
    return json_str

@app.post("/generate-embedding")
async def get_embeddings(request: Request):
    try:
        # Read the raw request
        body = await request.body()
        logger.info(f"Raw request size: {len(body)} bytes")
        
        try:
            # First try to parse the raw JSON
            try:
                data = json.loads(body)
            except json.JSONDecodeError:
                # If that fails, try to decode and clean the body
                body_str = body.decode('utf-8', errors='ignore')
                
                # Find all JSON objects in the string
                json_objects = []
                depth = 0
                start = -1
                
                for i, char in enumerate(body_str):
                    if char == '{':
                        if depth == 0:
                            start = i
                        depth += 1
                    elif char == '}':
                        depth -= 1
                        if depth == 0 and start != -1:
                            json_objects.append(body_str[start:i+1])
                            start = -1
                
                # If we found multiple JSON objects, take the first complete one
                if json_objects:
                    json_str = json_objects[0]
                    logger.info(f"Found {len(json_objects)} JSON objects, using first complete one")
                else:
                    raise ValueError("No complete JSON objects found in request body")
                
                # Clean the JSON string
                json_str = json_str.strip()
                # Remove any whitespace outside of quotes
                json_str = re.sub(r'\s+(?=(?:[^"]*"[^"]*")*[^"]*$)', '', json_str)
                
                try:
                    # Parse the cleaned JSON
                    data = json.loads(json_str)
                    # Validate it has the expected structure
                    if not isinstance(data, dict) or "chunks" not in data:
                        raise ValueError("Invalid JSON structure: missing 'chunks' key")
                except json.JSONDecodeError as je:
                    error_pos = je.pos
                    context = json_str[max(0, error_pos-50):min(len(json_str), error_pos+50)]
                    logger.error(f"JSON decode error at position {error_pos}: {context}")
                    raise
                
        except Exception as e:
            logger.error(f"JSON parsing error: {str(e)}")
            return JSONResponse(
                status_code=400,
                content={
                    "error": "Invalid JSON format",
                    "details": str(e)
                }
            )

        # Validate the chunks structure
        chunks = data.get("chunks", [])
        if not isinstance(chunks, list):
            return JSONResponse(
                status_code=400,
                content={"error": "Invalid chunks format - expected a list"}
            )

        logger.info(f"Processing {len(chunks)} chunks")
        results = []
        
        for i, chunk in enumerate(chunks):
            try:
                if not isinstance(chunk, dict):
                    logger.warning(f"Skipping chunk {i}: Not a dictionary")
                    continue
                    
                content = chunk.get("content")
                metadata = chunk.get("metadata", {})
                
                if not content:
                    logger.warning(f"Skipping chunk {i}: Missing content")
                    continue

                # Clean and normalize the content
                content = str(content).strip()
                # Remove any non-ASCII characters
                content = "".join(char for char in content if ord(char) < 128)
                
                if not content:
                    logger.warning(f"Skipping chunk {i}: Empty content after cleaning")
                    continue

                logger.info(f"Processing chunk {i} with content length: {len(content)}")
                embedding = await generate_embedding(content)
                
                if embedding is not None:
                    results.append({
                        "content": content,
                        "embedding": embedding.tolist(),
                        "metadata": metadata
                    })
                    logger.info(f"Successfully processed chunk {i}")
                else:
                    logger.warning(f"Failed to generate embedding for chunk {i}")
                    
            except Exception as chunk_error:
                logger.error(f"Error processing chunk {i}: {chunk_error}")
                continue

        if not results:
            return JSONResponse(
                status_code=500,
                content={"error": "Failed to generate any embeddings"}
            )

        # Create final response with strict JSON encoding
        response_data = {"results": results}
        response_json = json.dumps(response_data, ensure_ascii=True, separators=(',', ':'))
        
        # Validate the response JSON
        try:
            json.loads(response_json)
        except json.JSONDecodeError as je:
            logger.error(f"Invalid response JSON: {je}")
            return JSONResponse(
                status_code=500,
                content={"error": "Failed to generate valid response JSON"}
            )
        
        logger.info(f"Successfully generated embeddings for {len(results)} chunks")
        return JSONResponse(
            content=json.loads(response_json),
            media_type="application/json"
        )

    except Exception as e:
        logger.error(f"Embedding service failed: {str(e)}", exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"error": f"Internal server error: {str(e)}"}
        )

class ChunkItem(BaseModel):
    content: str
    embedding: List[float]
    metadata: Dict[str, Any]

class UploadRequest(BaseModel):
    results: List[ChunkItem]


import base64

def encode_web_url(web_url: str) -> str:
    """Encodes a web_url to a format suitable for Microsoft Graph /shares/u! endpoint."""
    encoded = base64.urlsafe_b64encode(web_url.encode("utf-8")).decode("utf-8").rstrip("=")
    return f"u!{encoded}"

@app.post("/upload-data")
async def upload_to_weaviate(req: UploadRequest, authorization: str = Header(None)):
    try:
        client = weaviate.connect_to_weaviate_cloud(
            cluster_url=WEAVIATE_URL,
            auth_credentials=AuthApiKey(WEAVIATE_API_KEY),
            skip_init_checks=True
        )
        # collection = client.collections.get("Chunks")
        collection = client.collections.get("Test")
        logger.info("Connected to Weaviate and accessed 'Test' collection")

        inserted = []
        # files_to_delete = set()  # No longer needed

        for i, chunk in enumerate(req.results):
            try:
                collection.data.insert(
                    properties={
                        "content": chunk.content,
                        "metadata": chunk.metadata
                    },
                    vector=np.array(chunk.embedding).tolist()
                )
                inserted.append(f"Inserted chunk {i} from {chunk.metadata.get('source_file')}")
                # files_to_delete.add(chunk.metadata.get("filepath"))  # Skip collecting for deletion
            except Exception as insert_error:
                logger.warning(f"Failed to insert chunk {i}: {insert_error}")

        # === Skip file deletion from OneDrive ===
        # for file_url in files_to_delete:
        #     if file_url:
        #         await delete_from_onedrive(file_url)

        return {"status": "success", "inserted": inserted}

    except Exception as e:
        logger.error(f"Weaviate upload error: {e}")
        return {"status": "error", "message": str(e)}
# ========== SCHEMA DEFINITIONS ==========
class SourceDocument(BaseModel):
    text: str = Field(..., description="The content of the source document")
    metadata: Dict[str, Any] = Field(..., description="Metadata associated with the document")
    relevance_score: float = Field(..., description="Relevance score of the document")

class CostDetails(BaseModel):
    input_tokens: int = Field(..., description="Number of input tokens used")
    output_tokens: int = Field(..., description="Number of output tokens generated")
    embedding_tokens: int = Field(..., description="Number of tokens used for embedding")
    costs: Dict[str, Any] = Field(..., description="Detailed cost breakdown")

class QueryRequest(BaseModel):
    query: str = Field(..., description="The query text to process")
    n_results: int = Field(3, description="Number of results to return", ge=1, le=10)
    model: str = Field("gpt-4o-mini", description="Model to use for completion")
    calculate_cost: bool = Field(True, description="Whether to include cost breakdown")

class QueryResponse(BaseModel):
    answer: str = Field(..., description="Generated answer to the query")
    sources: List[SourceDocument] = Field(..., description="List of source documents")
    cost_breakdown: Optional[CostDetails] = Field(None, description="Detailed cost information")

# ========== CONSTANTS ==========
MODEL_PRICING = {
    "gpt-4o-mini": {
        "input": 0.0005 / 1000,
        "output": 0.0015 / 1000,
        "embedding": 0.0001 / 1000
    },
    "gpt-4": {
        "input": 0.03 / 1000,
        "output": 0.06 / 1000,
        "embedding": 0.0001 / 1000
    }
}

# ========== HELPER FUNCTIONS ==========
def count_tokens(text: str, model: str = "gpt-4o-mini") -> int:
    """Count tokens in text using tiktoken for the specified model"""
    try:
        enc = tiktoken.encoding_for_model(model)
    except KeyError:
        enc = tiktoken.get_encoding("cl100k_base")
    return len(enc.encode(text))

def calculate_costs(
    prompt_tokens: int,
    completion_tokens: int,
    embedding_tokens: int,
    model: str
) -> Dict[str, Any]:
    """Calculate all costs in USD and INR"""
    USD_TO_INR = 83.50  # Current exchange rate
    
    model_price = MODEL_PRICING.get(model, MODEL_PRICING["gpt-4o-mini"])
    
    input_cost_usd = prompt_tokens * model_price["input"]
    output_cost_usd = completion_tokens * model_price["output"]
    embedding_cost_usd = embedding_tokens * model_price["embedding"]
    
    return {
        "input_tokens": prompt_tokens,
        "output_tokens": completion_tokens,
        "embedding_tokens": embedding_tokens,
        "input_cost_usd": input_cost_usd,
        "output_cost_usd": output_cost_usd,
        "embedding_cost_usd": embedding_cost_usd,
        "input_cost_inr": input_cost_usd * USD_TO_INR,
        "output_cost_inr": output_cost_usd * USD_TO_INR,
        "embedding_cost_inr": embedding_cost_usd * USD_TO_INR
    }

# ========== QUERY ENDPOINT ==========
@app.post("/query", response_model=QueryResponse)
async def process_query(request: QueryRequest):
    try:
        # Initialize all variables
        embedding_tokens = count_tokens(request.query)
        query_embedding = await generate_embedding1(request.query)
        
        if query_embedding is None:
            raise HTTPException(status_code=500, detail="Embedding generation failed")

        # Retrieve documents
        result = query_with_retries(query_embedding, request.n_results)
        if not hasattr(result, "objects"):
            raise HTTPException(status_code=500, detail="Invalid Weaviate response")

        documents = []
        metadatas = []
        distances = []
        
        for obj in result.objects:
            if hasattr(obj, "properties"):
                documents.append(obj.properties.get("content", ""))
                metadatas.append(obj.properties.get("metadata", {}))
                distances.append(getattr(obj.metadata, "distance", 0.0))

        if not documents:
            raise HTTPException(status_code=404, detail="No documents found")

        # Generate LLM response
        context = "\n\n".join(f"Document {i+1}: {doc}" for i, doc in enumerate(documents))
        prompt = f"Context:\n{context}\n\nQuestion: {request.query}\n\nAnswer:"
        
        llm_response, prompt_tokens, completion_tokens = await generate_completion(
            prompt=prompt,
            system_prompt = """
        You are a helpful assistant that answers questions using only the provided context.
Each answer must be well-structured, grammatically correct, and include relevant details such as source slide number and source file.
 
Your goal is to:
1. Extract meaningful, non-repetitive information from each context item.
2. Use commas, periods, and proper sentence structure.
3. Make sure the answer is insightful, with 2–3 key points per resource.
4. Avoid copying identical phrases across resources; ensure uniqueness.
5. If the context does not contain enough information, respond with: "I don't have enough information to answer this question based on the context."
 
Format the answer like this if the chunks are from ppt:
1. [Slide Number] Summary of the key insight with proper punctuation and grammar. Mention the unique takeaway.
2. [Slide Number] Different insight from another source. Ensure no repetition from the previous point.
3. [Slide Number] Another unique takeaway with source attribution.
 
Format the answer like this if the chunks are from docx:
1. Summary of the key insight with proper punctuation and grammar. Mention the unique takeaway.
2. Different insight from another source. Ensure no repetition from the previous point.
3. Another unique takeaway with source attribution.
 
If only 1 or 2 relevant contexts exist, only return those.
 
Only answer based on the provided context. Do not hallucinate or guess.
        """,
            model=request.model
        )

        # Build response
        response_data = {
            "answer": llm_response,
            "sources": [
                SourceDocument(
                    text=doc,
                    metadata=metadatas[i],
                    relevance_score=1.0 - distances[i]
                ) for i, doc in enumerate(documents)
            ]
        }

        # Add cost breakdown if requested
        if request.calculate_cost:
            cost_details = calculate_costs(
                prompt_tokens=prompt_tokens,
                completion_tokens=completion_tokens,
                embedding_tokens=embedding_tokens,
                model=request.model
            )
            
            response_data["cost_breakdown"] = CostDetails(
                input_tokens=cost_details["input_tokens"],
                output_tokens=cost_details["output_tokens"],
                embedding_tokens=cost_details["embedding_tokens"],
                costs={
                    "embedding": {
                        "usd": f"${cost_details['embedding_cost_usd']:.6f}",
                        "inr": f"₹{cost_details['embedding_cost_inr']:.4f}"
                    },
                    "summarization": {
                        "input": {
                            "usd": f"${cost_details['input_cost_usd']:.6f}",
                            "inr": f"₹{cost_details['input_cost_inr']:.4f}"
                        },
                        "output": {
                            "usd": f"${cost_details['output_cost_usd']:.6f}",
                            "inr": f"₹{cost_details['output_cost_inr']:.4f}"
                        }
                    },
                    "overall": {
                        "usd": f"${sum(cost_details[k] for k in ['input_cost_usd', 'output_cost_usd', 'embedding_cost_usd']):.6f}",
                        "inr": f"₹{sum(cost_details[k] for k in ['input_cost_inr', 'output_cost_inr', 'embedding_cost_inr']):.4f}"
                    }
                }
            )

        return response_data

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Query processing failed: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail="Internal server error")
####### Query Section ########
weaviate_client = weaviate.connect_to_weaviate_cloud(
    cluster_url=WEAVIATE_URL,
    auth_credentials=AuthApiKey(WEAVIATE_API_KEY),
    skip_init_checks=True,
)
# # Maximum retries for API calls
MAX_RETRIES = 3
RETRY_DELAY = 2
def query_with_retries(query_embedding, n_results):
    for attempt in range(MAX_RETRIES):
        try:
            # Initialize the Weaviate client
            client = weaviate.connect_to_weaviate_cloud(
                cluster_url=WEAVIATE_URL,
                auth_credentials=AuthApiKey(WEAVIATE_API_KEY),
                skip_init_checks=True  
            )
            
            # Access the collection
            collection = client.collections.get("Chunks")
            
            # Perform the query
            result = collection.query.near_vector(
                near_vector=query_embedding.tolist(),
                limit=n_results,
                return_metadata=MetadataQuery(distance=True)
            )
            return result
        
        except grpc._channel._InactiveRpcError:
            # Handle stale connection
            logger.warning("Stale connection detected. Reinitializing client...")
            time.sleep(RETRY_DELAY)  # Optional: Add a delay before retrying
            continue  # Retry the query
        
        except Exception as e:
            # Handle other exceptions
            logger.error(f"Query attempt {attempt + 1} failed: {str(e)}")
            if attempt < MAX_RETRIES - 1:
                time.sleep(RETRY_DELAY)  # Add a delay before retrying
                continue
            raise e
async def generate_embedding1(text: str) -> np.ndarray:
    """Generate embedding using Azure OpenAI"""
    try:
        async with ClientSession() as session:
            async with AiohttpTransport(client=session) as aiohttp_transport:
                httpx_client = openai.DefaultAsyncHttpxClient(transport=aiohttp_transport)
                client = AsyncAzureOpenAI(
                    azure_endpoint=AZURE_OPENAI_ENDPOINT,
                    api_key=AZURE_OPENAI_API_KEY,
                    api_version=AZURE_OPENAI_API_VERSION,
                    http_client=httpx_client
                )
                response = await client.embeddings.create(
                    model=EMBEDDING_MODEL,
                    input=[text],
                    encoding_format="float"
                )
                await httpx_client.aclose()
                return np.array(response.data[0].embedding)
    except Exception as e:
        logger.error(f"Embedding error: {e}")
        return None
async def generate_completion(prompt: str, system_prompt: str = None, model: str = "gpt-4o-mini"):
    try:
        async with ClientSession() as session:
            async with AiohttpTransport(client=session) as aiohttp_transport:
                httpx_client = openai.DefaultAsyncHttpxClient(transport=aiohttp_transport)
                client = AsyncAzureOpenAI(
                    azure_endpoint=AZURE_OPENAI_ENDPOINT,
                    api_key=AZURE_OPENAI_API_KEY,
                    api_version=AZURE_OPENAI_API_VERSION,
                    http_client=httpx_client
                )
                try:
                    messages = []
                    if system_prompt:
                        messages.append({"role": "system", "content": system_prompt})
                    messages.append({"role": "user", "content": prompt})

                    logger.info(f"Generating completion with model: {model}")
                    response = await client.chat.completions.create(
                        model=model,
                        messages=messages,
                        temperature=0.3,
                        max_tokens=512
                    )

                    if not response or not response.choices:
                        logger.error("Empty response from OpenAI completion API")
                        return "", 0, 0

                    usage = getattr(response, "usage", None)
                    prompt_tokens = usage.prompt_tokens if usage else 0
                    completion_tokens = usage.completion_tokens if usage else 0

                    return response.choices[0].message.content.strip(), prompt_tokens, completion_tokens
                finally:
                    await httpx_client.aclose()
    except Exception as e:
        logger.error(f"Error generating completion: {str(e)}")
        return "", 0, 0
 



# Add a function to calculate costs
MODEL_PRICING = {
    "gpt-4o-mini": {
        "input": 0.15 / 1_000_000,   # $0.15 per million input tokens
        "output": 0.60 / 1_000_000,  # $0.60 per million output tokens
        "embedding": 0.10 / 1_000_000  # $0.10 per million embedding tokens
    }
}

def calculate_costs(prompt_tokens: int, completion_tokens: int, embedding_tokens: int = 0):
    """Calculate costs based on your specified rates"""
    USD_TO_INR = 86.38
    
    # Calculate USD costs
    input_cost_usd = prompt_tokens * MODEL_PRICING["gpt-4o-mini"]["input"]
    output_cost_usd = completion_tokens * MODEL_PRICING["gpt-4o-mini"]["output"]
    embedding_cost_usd = embedding_tokens * MODEL_PRICING["gpt-4o-mini"]["embedding"]
    
    # Convert to INR
    input_cost_inr = input_cost_usd * USD_TO_INR
    output_cost_inr = output_cost_usd * USD_TO_INR
    embedding_cost_inr = embedding_cost_usd * USD_TO_INR
    
    return {
        "input_tokens": prompt_tokens,
        "output_tokens": completion_tokens,
        "embedding_tokens": embedding_tokens,
        "costs_usd": {
            "input": input_cost_usd,
            "output": output_cost_usd,
            "embedding": embedding_cost_usd,
            "total": input_cost_usd + output_cost_usd + embedding_cost_usd
        },
        "costs_inr": {
            "input": input_cost_inr,
            "output": output_cost_inr,
            "embedding": embedding_cost_inr,
            "total": input_cost_inr + output_cost_inr + embedding_cost_inr
        }
    }

def count_tokens(text: str, model: str = "gpt-4o-mini") -> int:
    """Count tokens in text using tiktoken for the specified model"""
    try:
        enc = tiktoken.encoding_for_model(model)
    except KeyError:
        enc = tiktoken.get_encoding("cl100k_base")
    return len(enc.encode(text))

def calculate_costs(
    prompt_tokens: int, 
    completion_tokens: int, 
    embedding_tokens: int = 0, 
    model: str = "gpt-4o-mini"
) -> dict:
    """Calculate all costs in USD and INR"""
    USD_TO_INR = 86.43  # Update to current exchange rate
    
    # Get pricing for the specified model, default to gpt-4o-mini if not found
    model_price = MODEL_PRICING.get(model, MODEL_PRICING["gpt-4o-mini"])
    
    # Calculate individual costs
    input_cost_usd = prompt_tokens * model_price["input"]
    output_cost_usd = completion_tokens * model_price["output"]
    embedding_cost_usd = embedding_tokens * model_price["embedding"]
    
    # Convert to INR
    input_cost_inr = input_cost_usd * USD_TO_INR
    output_cost_inr = output_cost_usd * USD_TO_INR
    embedding_cost_inr = embedding_cost_usd * USD_TO_INR
    
    # Calculate totals
    total_cost_usd = input_cost_usd + output_cost_usd + embedding_cost_usd
    total_cost_inr = input_cost_inr + output_cost_inr + embedding_cost_inr
    
    return {
        "input_tokens": prompt_tokens,
        "output_tokens": completion_tokens,
        "embedding_tokens": embedding_tokens,
        "input_cost_usd": input_cost_usd,
        "output_cost_usd": output_cost_usd,
        "embedding_cost_usd": embedding_cost_usd,
        "total_cost_usd": total_cost_usd,
        "input_cost_inr": input_cost_inr,
        "output_cost_inr": output_cost_inr,
        "embedding_cost_inr": embedding_cost_inr,
        "total_cost_inr": total_cost_inr
    }

# Update your query endpoint to use these functions
@app.post("/query", response_model=QueryResponse)
async def process_query(request: QueryRequest):
    try:
        logger.info(f"Processing query: {request.query}")
        
        # Initialize all variables we'll use
        embedding_tokens = 0
        prompt_tokens = 0
        completion_tokens = 0
        llm_response = ""
        documents = []
        metadatas = []
        distances = []
        
        # Step 1: Generate embedding and count tokens
        embedding_tokens = count_tokens(request.query)
        query_embedding = await generate_embedding1(request.query)
        if query_embedding is None:
            raise HTTPException(status_code=500, detail="Failed to generate query embedding")
        
        # Step 2: Retrieve documents from Weaviate
        result = query_with_retries(query_embedding, request.n_results)
        
        if not hasattr(result, "objects") or not isinstance(result.objects, list):
            raise HTTPException(status_code=500, detail="Invalid response format from Weaviate")

        # Process retrieved documents
        for obj in result.objects:
            if hasattr(obj, "properties"):
                documents.append(obj.properties.get("content", ""))
                metadatas.append(obj.properties.get("metadata", {}))
                distances.append(obj.metadata.distance if hasattr(obj.metadata, "distance") else 0.0)

        if not documents:
            raise HTTPException(status_code=404, detail="No relevant documents found")

        # Prepare context for LLM
        context = "\n\n".join([f"Document {i+1}: {doc}" for i, doc in enumerate(documents)])
        
        # Step 3: Generate answer using LLM
        system_prompt = """..."""  # Your existing system prompt here

        prompt = f"""Context:
{context}

Question: {request.query}

Please answer the question based only on the provided context."""

        # Generate completion and get token counts
        llm_response, prompt_tokens, completion_tokens = await generate_completion(
            prompt=prompt,
            system_prompt=system_prompt,
            model=request.model
        )

        # Prepare response data structure
        response_data = {
            "answer": llm_response,
            "sources": [{
                "text": doc,
                "metadata": metadatas[i] if i < len(metadatas) else {},
                "relevance_score": 1.0 - distances[i] if i < len(distances) else 0
            } for i, doc in enumerate(documents)]
        }

        # Calculate costs if requested
        if request.calculate_cost:
            cost_details = calculate_costs(
                prompt_tokens=prompt_tokens,
                completion_tokens=completion_tokens,
                embedding_tokens=embedding_tokens,
                model=request.model
            )
            
            response_data["cost_breakdown"] = {
                "input_tokens": cost_details["input_tokens"],
                "output_tokens": cost_details["output_tokens"],
                "embedding_tokens": cost_details["embedding_tokens"],
                "costs": {
                    "embedding": {
                        "usd": f"${cost_details['embedding_cost_usd']:.6f}",
                        "inr": f"₹{cost_details['embedding_cost_inr']:.4f}"
                    },
                    "summarization": {
                        "input": {
                            "usd": f"${cost_details['input_cost_usd']:.6f}",
                            "inr": f"₹{cost_details['input_cost_inr']:.4f}"
                        },
                        "output": {
                            "usd": f"${cost_details['output_cost_usd']:.6f}",
                            "inr": f"₹{cost_details['output_cost_inr']:.4f}"
                        },
                        "total": {
                            "usd": f"${(cost_details['input_cost_usd'] + cost_details['output_cost_usd']):.6f}",
                            "inr": f"₹{(cost_details['input_cost_inr'] + cost_details['output_cost_inr']):.4f}"
                        }
                    },
                    "overall": {
                        "usd": f"${cost_details['total_cost_usd']:.6f}",
                        "inr": f"₹{cost_details['total_cost_inr']:.4f}"
                    }
                }
            }

        return response_data

    except Exception as e:
        logger.error(f"Error processing query: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))
 
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)