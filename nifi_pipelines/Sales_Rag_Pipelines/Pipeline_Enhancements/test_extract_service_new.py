import os
import sys
import json
import logging
from io import BytesIO
from typing import List, Dict, Any, Tuple
from pptx import Presentation
from pptx.shapes.base import BaseShape
from docx import Document
from dotenv import load_dotenv
import numpy as np
from azure.storage.blob import BlobServiceClient, BlobClient, ContentSettings
import re
import subprocess
import uuid
import hashlib
import zipfile
import tempfile

# Load environment variables
load_dotenv()

# === Configuration ===
AZURE_ACCOUNT_URL = os.environ["AZURE_BLOB_ACCOUNT_URL"]
AZURE_CONTAINER_NAME = os.environ.get("AZURE_CONTAINER_NAME", "ppt-slide-storage")
AZURE_SAS_TOKEN = os.environ["AZURE_SAS_TOKEN"]

# === Logging ===
logging.basicConfig(
    level=logging.INFO, 
    format='%(asctime)s - %(levelname)s - %(message)s',
    stream=sys.stderr
)
logger = logging.getLogger(__name__)

# Create service and container clients
blob_service_client = BlobServiceClient(account_url=AZURE_ACCOUNT_URL, credential=AZURE_SAS_TOKEN)

def get_blob_client(blob_name: str) -> BlobClient:
    container_client = blob_service_client.get_container_client(AZURE_CONTAINER_NAME)
    return container_client.get_blob_client(blob_name)

def blob_exists(blob_name: str) -> bool:
    blob_client = get_blob_client(blob_name)
    try:
        blob_client.get_blob_properties()
        return True
    except Exception:
        return False

def download_json_from_blob(blob_name: str) -> dict:
    blob_client = get_blob_client(blob_name)
    stream = blob_client.download_blob()
    return json.loads(stream.readall())

def upload_json_to_blob(blob_name: str, data: dict):
    blob_client = get_blob_client(blob_name)
    blob_client.upload_blob(
        json.dumps(data),
        overwrite=True,
        content_settings=ContentSettings(content_type="application/json")
    )

def hash_pptx_file(file_bytes: BytesIO) -> str:
    file_bytes.seek(0)
    hasher = hashlib.sha256()
    while chunk := file_bytes.read(8192):
        hasher.update(chunk)
    file_bytes.seek(0)
    return hasher.hexdigest()

def extract_text_from_docx(file_bytes: BytesIO) -> List[str]:
    doc = Document(file_bytes)
    chunks = []
    current_chunk = []
    current_size = 0
    max_chunk_size = 1028
    min_chunk_size = 200

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if current_size + len(text) > max_chunk_size and current_chunk:
            chunk_text = " ".join(current_chunk)
            if len(chunk_text) >= min_chunk_size:
                chunks.append(chunk_text)
            current_chunk = []
            current_size = 0

        current_chunk.append(text)
        current_size += len(text)

    if current_chunk:
        chunk_text = " ".join(current_chunk)
        if len(chunk_text) >= min_chunk_size:
            chunks.append(chunk_text)

    return chunks

def clean_extracted_text(raw_text: str) -> str:
    """Enhanced text cleaning function"""
    text = " ".join(raw_text.split())
    
    def fix_spaced_characters(match):
        chars = match.group(0)
        return re.sub(r'(?<=\w)\s+(?=\w)', '', chars)
    
    text = re.sub(r'\b\w(?:\s+\w){2,}\b', fix_spaced_characters, text)
    text = re.sub(r'\b([a-z])\s+([a-z])\s+([a-z])', r'\1\2\3', text)
    text = re.sub(r'\b([A-Z])\s+([a-z])\s+([a-z])', r'\1\2\3', text)
    text = re.sub(r'\b([a-z])\s+([a-z])\b(?=\s+[a-z]\s+[a-z])', r'\1\2', text)
    
    # Fix common broken patterns
    text = re.sub(r'\bFo\s+ra\b', 'For a', text)
    text = re.sub(r'\bC\s+us\s+to\s+me\b', 'Customer', text)
    text = re.sub(r'\bE\s+xp\s+er\s+ie\s+nc\s+e\b', 'Experience', text)
    text = re.sub(r'\bP\s+la\s+tf\s+or\s+m\b', 'Platform', text)
    
    common_broken_patterns = {
        r'\bGen\s+AI\b': 'GenAI',
        r'\bTrans\s+for\s+ma\s+tion\b': 'Transformation',
        r'\bChal\s+len\s+ges\b': 'Challenges',
        r'\bSo\s+lu\s+tions\b': 'Solutions',
        r'\bOut\s+co\s+mes\b': 'Outcomes',
    }
    
    for pattern, replacement in common_broken_patterns.items():
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
    
    def merge_likely_word(match):
        chars = match.group(0)
        merged = re.sub(r'\s+', '', chars)
        if len(merged) >= 3 and re.search(r'[aeiouAEIOU]', merged):
            return merged
        return chars
    
    text = re.sub(r'\b[A-Za-z](?:\s+[A-Za-z]){2,}\b', merge_likely_word, text)
    text = re.sub(r'\s{2,}', ' ', text)
    text = re.sub(r'([a-z])([A-Z][a-z])', r'\1 \2', text)
    
    return text.strip()

def extract_structured_text_from_pptx_slide(slide) -> Dict[str, Any]:
    """
    Extract text from a single slide maintaining structure based on position and content
    """
    text_elements = []
    
    # Extract all text shapes with their positions and properties
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # Get position and size
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            
            # Get text content
            text_content = shape.text.strip()
            
            # Try to determine text type based on formatting and position
            text_type = classify_text_element(shape, text_content, top, left, width, height)
            
            text_elements.append({
                'text': text_content,
                'type': text_type,
                'position': {'left': left, 'top': top, 'width': width, 'height': height},
                'font_size': get_average_font_size(shape),
                'is_bold': is_text_bold(shape)
            })
    
    # Sort elements by logical reading order
    structured_elements = organize_text_elements(text_elements)
    
    return structured_elements

def classify_text_element(shape, text_content: str, top: int, left: int, width: int, height: int) -> str:
    """
    Classify text elements based on content, position, and formatting - now handles ANY section headers
    """
    text_lower = text_content.lower().strip()
    text_clean = text_content.strip()
    
    # Check if this looks like a main title (usually at the top, longer text)
    if top < 100000 and len(text_clean) > 30:  # Very top of slide
        return 'main_title'
    
    # Check if this looks like a subtitle (usually under main title, starts with "for")
    if text_lower.startswith('for a') or text_lower.startswith('for the') or text_lower.startswith('for an'):
        return 'subtitle'
    
    # Check if this is a section header based on characteristics:
    # 1. Short text (typically 1-4 words)
    # 2. Position-based (not at very top or very bottom)
    # 3. Likely to be bold or larger font
    # 4. Doesn't start with bullet points or numbers
    
    word_count = len(text_clean.split())
    is_likely_header = (
        word_count <= 4 and  # Short text
        len(text_clean) < 50 and  # Not too long
        100000 < top < 500000 and  # Middle area of slide (not title, not bottom)
        not text_clean.startswith(('•', '-', '*', '1.', '2.', '3.', '4.', '5.')) and  # Not a bullet/list
        not re.match(r'^\d+[.\):]', text_clean) and  # Not numbered list
        not re.search(r'\d{1,3}%', text_clean) and  # Not a percentage/metric
        text_clean[0].isupper()  # Starts with capital letter
    )
    
    if is_likely_header:
        # Create a generic section header
        section_name = text_clean.lower().replace(' ', '_').replace('-', '_')
        # Clean up the section name
        section_name = re.sub(r'[^\w_]', '', section_name)
        return f'section_header:{section_name}'
    
    # Check for bullet points or list items
    if (text_clean.startswith('•') or text_clean.startswith('-') or text_clean.startswith('*') or
        re.match(r'^\d+[.\):]', text_clean)):
        return 'bullet_point'
    
    # Check for percentage or metrics (standalone numbers with %)
    if re.search(r'\b\d{1,3}%\b', text_clean) and len(text_clean) < 20:
        return 'metric'
    
    # Check for very top area headers (company names, etc.)
    if top < 50000:
        return 'top_header'
    
    # Check for bottom area text (often company info, page numbers)
    if top > 600000:  # Bottom area
        return 'footer_text'
    
    # Default classification based on length and position
    if len(text_clean) < 30:
        return 'short_text'
    else:
        return 'body_text'

def get_average_font_size(shape) -> float:
    """Get average font size from a text shape"""
    try:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        return float(run.font.size.pt)
        return 12.0  # Default
    except:
        return 12.0

def is_text_bold(shape) -> bool:
    """Check if text is bold"""
    try:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.bold:
                        return True
        return False
    except:
        return False

def organize_text_elements(text_elements: List[Dict]) -> Dict[str, Any]:
    """
    Organize text elements into logical structure - now handles ANY section headers dynamically
    """
    # Sort by vertical position first (top to bottom), then horizontal (left to right)
    text_elements.sort(key=lambda x: (x['position']['top'], x['position']['left']))
    
    structured_content = {
        'main_title': '',
        'subtitle': '',
        'top_headers': [],  # Company names, logos, etc.
        'sections': {},
        'footer_text': []   # Bottom text, page numbers, etc.
    }
    
    current_section = None
    current_section_content = []
    detected_sections = []  # Keep track of section order
    
    for element in text_elements:
        text_type = element['type']
        text_content = element['text'].strip()
        
        if text_type == 'main_title':
            structured_content['main_title'] = text_content
        elif text_type == 'subtitle':
            structured_content['subtitle'] = text_content
        elif text_type == 'top_header':
            structured_content['top_headers'].append(text_content)
        elif text_type == 'footer_text':
            structured_content['footer_text'].append(text_content)
        elif text_type.startswith('section_header:'):
            # Dynamic section header handling
            # Save previous section if exists
            if current_section and current_section_content:
                if current_section not in structured_content['sections']:
                    structured_content['sections'][current_section] = []
                structured_content['sections'][current_section].extend(current_section_content)
            
            # Extract section name from the type
            section_name = text_type.split(':', 1)[1]  # Get part after 'section_header:'
            current_section = section_name
            current_section_content = []
            
            # Track the order of sections as they appear
            if section_name not in detected_sections:
                detected_sections.append(section_name)
                
            # Also store the header text itself
            current_section_content.append({
                'type': 'section_title',
                'text': text_content,
                'position': element['position']
            })
            
        else:
            # Add to current section or create general section
            content_item = {
                'type': text_type,
                'text': text_content,
                'position': element['position']
            }
            
            if current_section:
                current_section_content.append(content_item)
            else:
                # No section header found yet, add to general content
                if 'general' not in structured_content['sections']:
                    structured_content['sections']['general'] = []
                    if 'general' not in detected_sections:
                        detected_sections.append('general')
                structured_content['sections']['general'].append(content_item)
    
    # Don't forget the last section
    if current_section and current_section_content:
        if current_section not in structured_content['sections']:
            structured_content['sections'][current_section] = []
        structured_content['sections'][current_section].extend(current_section_content)
    
    # Store the order of sections for later use
    structured_content['section_order'] = detected_sections
    
    return structured_content

def create_structured_chunks(structured_content: Dict[str, Any], slide_number: int, slide_url: str) -> List[str]:
    """
    Create chunks that maintain logical structure - now handles ANY sections dynamically
    """
    chunks = []
    
    # Build the structured text
    slide_text_parts = []
    
    # Add top headers (company names, etc.)
    if structured_content.get('top_headers'):
        slide_text_parts.extend(structured_content['top_headers'])
    
    # Add title and subtitle
    if structured_content['main_title']:
        slide_text_parts.append(structured_content['main_title'])
    if structured_content['subtitle']:
        slide_text_parts.append(structured_content['subtitle'])
    
    # Add sections in the order they appeared on the slide
    section_order = structured_content.get('section_order', [])
    
    for section_name in section_order:
        if section_name in structured_content['sections']:
            section_content = structured_content['sections'][section_name]
            
            # Find the section title from the content
            section_title = None
            section_items = []
            
            for item in section_content:
                if item['type'] == 'section_title':
                    section_title = item['text']
                else:
                    section_items.append(item['text'])
            
            # Add section with its title
            if section_title:
                slide_text_parts.append(f"\n{section_title}:")
            elif section_name != 'general':
                # If no title found, create one from section name
                display_name = section_name.replace('_', ' ').title()
                slide_text_parts.append(f"\n{display_name}:")
            
            # Add section content
            slide_text_parts.extend(section_items)
    
    # Add footer text if any
    if structured_content.get('footer_text'):
        slide_text_parts.extend(structured_content['footer_text'])
    
    # Join all parts
    complete_slide_text = " ".join(slide_text_parts)
    complete_slide_text = clean_extracted_text(complete_slide_text)
    
    # Create the final slide content
    slide_content = f"[Slide {slide_number}] {complete_slide_text}"
    
    # For now, return as single chunk (you can split later if needed)
    chunks.append({
        "content": slide_content,
        "slide_number": slide_number,
        "slide_url": slide_url
        # Removed structure from here since we don't want it in the final output
    })
    
    return chunks

def extract_text_from_pptx_structured(file_bytes: BytesIO, web_url: str, server_base_url: str = None, original_filename: str = None) -> List[Dict]:
    """
    Extract text from PowerPoint maintaining logical structure
    """
    chunks = []
    
    # First, handle the slide URL generation (same as before)
    pptx_hash = hash_pptx_file(file_bytes)
    metadata_blob_name = f"slide_metadata/{pptx_hash}.json"

    if blob_exists(metadata_blob_name):
        logger.info(f"Found existing metadata for hash {pptx_hash}")
        slide_metadata = download_json_from_blob(metadata_blob_name)
        def extract_slide_number_from_url(url):
            match = re.search(r"slide_(\d+)\.pdf", url)
            return int(match.group(1)) if match else 0
        slide_urls = sorted(
            slide_metadata["slide_urls"],
            key=extract_slide_number_from_url
        )
    else:
        logger.info(f"No metadata found, processing {original_filename}")
        slide_urls = []
        session_id = str(uuid.uuid4())

        with tempfile.TemporaryDirectory() as tmpdir:
            input_pptx_path = os.path.join(tmpdir, original_filename)
            with open(input_pptx_path, "wb") as f:
                f.write(file_bytes.getvalue())

            full_pdf_path = convert_pptx_to_pdf(input_pptx_path, tmpdir)
            output_pattern = os.path.join(tmpdir, "slide_%d.pdf")
            split_pdf_to_slides(full_pdf_path, output_pattern)

            slide_pdfs = sorted(
                [f for f in os.listdir(tmpdir) if re.match(r"slide_\d+\.pdf", f)],
                key=lambda x: int(re.search(r"slide_(\d+)\.pdf", x).group(1))
            )
            for idx, pdf_name in enumerate(slide_pdfs):
                src_pdf = os.path.join(tmpdir, pdf_name)
                blob_key = f"slides/{session_id}/{pdf_name}"

                with open(src_pdf, "rb") as data:
                    blob_client = get_blob_client(blob_key)
                    blob_client.upload_blob(
                        data,
                        overwrite=True,
                        content_settings=ContentSettings(content_type="application/pdf")
                    )

                slide_urls.append(blob_client.url)

        upload_json_to_blob(metadata_blob_name, {"slide_urls": slide_urls})
        logger.info(f"Saved metadata JSON for {pptx_hash}")

    # Now extract structured text directly from PowerPoint
    file_bytes.seek(0)
    presentation = Presentation(file_bytes)
    
    for slide_idx, slide in enumerate(presentation.slides):
        slide_number = slide_idx + 1
        # Fixed slide URL indexing - use slide_idx (0-based) to match the slide_urls array
        slide_url = slide_urls[slide_idx] if slide_idx < len(slide_urls) else ""
        
        try:
            # Extract structured content from this slide
            structured_content = extract_structured_text_from_pptx_slide(slide)
            
            # Create structured chunks
            slide_chunks = create_structured_chunks(structured_content, slide_number, slide_url)
            chunks.extend(slide_chunks)
            
        except Exception as e:
            logger.warning(f"Failed to extract structured text from slide {slide_number}: {e}")
            # Fallback to simple text extraction
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += " " + shape.text
            
            if slide_text.strip():
                slide_content = f"[Slide {slide_number}] {clean_extracted_text(slide_text)}"
                chunks.append({
                    "content": slide_content,
                    "slide_number": slide_number,
                    "slide_url": slide_url
                })

    return chunks

def convert_pptx_to_pdf(input_pptx_path: str, output_dir: str) -> str:
    existing_pdfs = set(f for f in os.listdir(output_dir) if f.lower().endswith(".pdf"))
    subprocess.run([
        "soffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, input_pptx_path
    ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    new_pdfs = [f for f in os.listdir(output_dir) if f.lower().endswith(".pdf") and f not in existing_pdfs]
    if not new_pdfs:
        raise Exception("LibreOffice conversion failed: No PDF created.")
    return os.path.join(output_dir, new_pdfs[0])

def split_pdf_to_slides(input_pdf_path: str, output_pattern: str):
    subprocess.run(["pdfseparate", input_pdf_path, output_pattern], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

def detect_file_type(binary_data: bytes) -> str:
    if binary_data.startswith(b'PK\x03\x04'):
        try:
            with zipfile.ZipFile(BytesIO(binary_data), 'r') as zip_file:
                file_list = zip_file.namelist()
                if 'word/document.xml' in file_list:
                    return 'docx'
                elif 'ppt/presentation.xml' in file_list:
                    return 'pptx'
        except:
            pass
    if binary_data.startswith(b'%PDF'):
        return 'pdf'
    return 'unknown'

def parse_arguments():
    filename = None
    web_url = None
    
    if len(sys.argv) >= 3:
        filename = sys.argv[1]
        web_url = sys.argv[2]
        logger.info(f"Parsed from separate args: filename={filename}, web_url={web_url}")
    elif len(sys.argv) == 2:
        arg_string = sys.argv[1]
        if ';' in arg_string:
            parts = arg_string.split(';')
            if len(parts) >= 2:
                filename = parts[0]
                web_url = parts[1]
                logger.info(f"Parsed from semicolon-separated args: filename={filename}, web_url={web_url}")
        else:
            filename = arg_string
            web_url = arg_string
            logger.info(f"Single argument provided: {filename}")
    
    if not filename:
        filename = os.environ.get('filename', os.environ.get('onedrive_filename', 'unknown_file'))
        logger.info(f"Got filename from environment: {filename}")
    
    if not web_url:
        web_url = os.environ.get('web_url', filename)
        logger.info(f"Got web_url from environment or defaulted to filename: {web_url}")
    
    return filename, web_url

def main():
    try:
        filename, web_url = parse_arguments()
        binary_data = sys.stdin.buffer.read()
        
        if not binary_data:
            raise ValueError("No input data received")
        
        if '.' not in filename:
            detected_type = detect_file_type(binary_data)
            if detected_type != 'unknown':
                filename = f"{filename}.{detected_type}"
        
        file_bytes = BytesIO(binary_data)
        filename_lower = filename.lower()
        
        if filename_lower.endswith(".docx"):
            chunks = extract_text_from_docx(file_bytes)
            file_type = "docx"
        elif filename_lower.endswith(".pptx"):
            chunks = extract_text_from_pptx_structured(file_bytes, web_url, original_filename=filename)
            file_type = "pptx"
        else:
            detected_type = detect_file_type(binary_data)
            if detected_type == 'docx':
                chunks = extract_text_from_docx(file_bytes)
                file_type = "docx"
            elif detected_type == 'pptx':
                chunks = extract_text_from_pptx_structured(file_bytes, web_url, original_filename=filename)
                file_type = "pptx"
            else:
                raise ValueError(f"Unsupported file type: {filename}")

        chunk_output = []
        for idx, chunk in enumerate(chunks):
            if isinstance(chunk, dict):
                text = chunk.get("content", "")
                slide_number = chunk.get("slide_number")
                slide_url = chunk.get("slide_url")
            else:
                text = str(chunk)
                slide_number = None
                slide_url = None

            cleaned_text = " ".join(text.split())
            cleaned_text = "".join(char for char in cleaned_text if ord(char) < 128)

            # Create simplified metadata with only the requested fields
            metadata = {
                "source_file": filename,
                "filepath": web_url,
                "chunk_index": idx,
                "file_type": file_type
            }
            
            # Only add slide-specific fields for PowerPoint files when they exist
            if slide_number is not None:
                metadata["slide_number"] = slide_number
            if slide_url is not None:
                metadata["slide_url"] = slide_url

            chunk_output.append({"content": cleaned_text, "metadata": metadata})

        result = {"chunks": chunk_output}
        json_output = json.dumps(result, ensure_ascii=True)
        
        sys.stdout.flush()
        print(json_output, flush=True)
        
        logger.info(f"Successfully processed {len(chunk_output)} structured chunks from {filename}")

    except Exception as e:
        logger.error(f"Text extraction failed: {str(e)}")
        error_result = {"error": str(e), "status": "failed"}
        print(json.dumps(error_result))
        sys.exit(1)

if __name__ == "__main__":
    main()