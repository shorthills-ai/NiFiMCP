import sys
import os
import shutil
import tempfile
import zipfile
import xml.etree.ElementTree as ET
import re
from io import BytesIO
from pptx import Presentation

# ############################################################################
# IMPORTANT:
# 1. This script requires the 'python-pptx' library to be installed in the
#    environment where NiFi's ExecuteScript processor runs.
#
# 2. You MUST configure the path to your local PowerPoint file below.
# ############################################################################

# --- CONFIGURATION ---
LOCAL_PPTX_TO_APPEND = "/path/to/your/local_presentation.pptx"
# --- END CONFIGURATION ---


# --- START: High-fidelity slide copy functions (adapted from your provided script) ---

def copy_single_slide(input_file, output_file, slide_number):
    """
    Copy a specific slide from a PPTX file to create a new PPTX file (single-slide, valid for PowerPoint),
    and copy all referenced media files.
    """
    if not os.path.exists(input_file):
        raise FileNotFoundError(f"Input file '{input_file}' not found.")
    if not input_file.lower().endswith('.pptx'):
        raise ValueError("Input file must be a .pptx file")
    with tempfile.TemporaryDirectory() as temp_dir:
        extract_path = os.path.join(temp_dir, "extracted")
        output_path = os.path.join(temp_dir, "output")
        with zipfile.ZipFile(input_file, 'r') as zip_ref:
            zip_ref.extractall(extract_path)
        slides_dir = os.path.join(extract_path, "ppt", "slides")
        slide_files = [f for f in os.listdir(slides_dir) if f.startswith('slide') and f.endswith('.xml')]
        slide_files.sort(key=lambda x: int(re.search(r'slide(\d+)\.xml', x).group(1)))
        if slide_number < 1 or slide_number > len(slide_files):
            raise ValueError(f"Slide number {slide_number} not found. Available slides: 1-{len(slide_files)}")
        target_slide = slide_files[slide_number - 1]
        os.makedirs(output_path, exist_ok=True)
        copy_complete_structure(extract_path, output_path)
        clean_slides_directory(output_path, target_slide)
        slides_rels_dir = os.path.join(extract_path, "ppt", "slides", "_rels")
        target_slide_num = re.search(r'slide(\d+)\.xml', target_slide).group(1)
        rels_file = os.path.join(slides_rels_dir, f'slide{target_slide_num}.xml.rels')
        media_files = set()
        if os.path.exists(rels_file):
            tree = ET.parse(rels_file)
            root = tree.getroot()
            for rel in root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                target = rel.get('Target')
                if target and (target.startswith('../media/') or target.startswith('media/')):
                    media_name = os.path.basename(target)
                    media_files.add(media_name)
        src_media_dir = os.path.join(extract_path, "ppt", "media")
        dst_media_dir = os.path.join(output_path, "ppt", "media")
        if media_files and os.path.exists(src_media_dir):
            os.makedirs(dst_media_dir, exist_ok=True)
            for media_name in media_files:
                src_media = os.path.join(src_media_dir, media_name)
                dst_media = os.path.join(dst_media_dir, media_name)
                if os.path.exists(src_media):
                    shutil.copy2(src_media, dst_media)
        update_presentation_structure(output_path)
        with zipfile.ZipFile(output_file, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
            for root, dirs, files in os.walk(output_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, output_path)
                    zip_ref.write(file_path, arcname)

def copy_complete_structure(extract_path, output_path):
    shutil.copytree(extract_path, output_path, dirs_exist_ok=True)

def clean_slides_directory(output_path, target_slide):
    slides_dir = os.path.join(output_path, "ppt", "slides")
    slides_rels_dir = os.path.join(output_path, "ppt", "slides", "_rels")
    target_slide_num = re.search(r'slide(\d+)\.xml', target_slide).group(1)
    for file in os.listdir(slides_dir):
        if file.startswith('slide') and file.endswith('.xml') and file != target_slide:
            os.remove(os.path.join(slides_dir, file))
    if os.path.exists(slides_rels_dir):
        for file in os.listdir(slides_rels_dir):
            if file.startswith('slide') and file.endswith('.xml.rels'):
                slide_num = re.search(r'slide(\d+)\.xml\.rels', file).group(1)
                if slide_num != target_slide_num:
                    os.remove(os.path.join(slides_rels_dir, file))
    if target_slide != 'slide1.xml':
        old_slide_path = os.path.join(slides_dir, target_slide)
        new_slide_path = os.path.join(slides_dir, 'slide1.xml')
        os.rename(old_slide_path, new_slide_path)
        old_rel_path = os.path.join(slides_rels_dir, f'slide{target_slide_num}.xml.rels')
        new_rel_path = os.path.join(slides_rels_dir, 'slide1.xml.rels')
        if os.path.exists(old_rel_path):
            os.rename(old_rel_path, new_rel_path)

def update_presentation_structure(output_path):
    pres_file = os.path.join(output_path, "ppt", "presentation.xml")
    if os.path.exists(pres_file):
        ET.register_namespace('', 'http://schemas.openxmlformats.org/presentationml/2006/main')
        ET.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')
        ET.register_namespace('r', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships')
        tree = ET.parse(pres_file)
        root = tree.getroot()
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        slide_id_list = root.find('.//p:sldIdLst', ns)
        if slide_id_list is not None:
            slide_id_list.clear()
            slide_id = ET.SubElement(slide_id_list, '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
            slide_id.set('id', '256')
            slide_id.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', 'rId2')
        tree.write(pres_file, encoding='utf-8', xml_declaration=True)
    pres_rels_file = os.path.join(output_path, "ppt", "_rels", "presentation.xml.rels")
    if os.path.exists(pres_rels_file):
        tree = ET.parse(pres_rels_file)
        root = tree.getroot()
        relationships_to_remove = [rel for rel in root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship') if rel.get('Type') == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide']
        for rel in relationships_to_remove:
            root.remove(rel)
        new_rel = ET.SubElement(root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
        new_rel.set('Id', 'rId2')
        new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        new_rel.set('Target', 'slides/slide1.xml')
        tree.write(pres_rels_file, encoding='utf-8', xml_declaration=True)
    app_file = os.path.join(output_path, "docProps", "app.xml")
    if os.path.exists(app_file):
        tree = ET.parse(app_file)
        root = tree.getroot()
        slides_elem = root.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/extended-properties}Slides')
        if slides_elem is not None: slides_elem.text = '1'
        pages_elem = root.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/extended-properties}Pages')
        if pages_elem is not None: pages_elem.text = '1'
        tree.write(app_file, encoding='utf-8', xml_declaration=True)
    content_types_file = os.path.join(output_path, "[Content_Types].xml")
    if os.path.exists(content_types_file):
        tree = ET.parse(content_types_file)
        root = tree.getroot()
        overrides_to_remove = [override for override in root.findall('.//{http://schemas.openxmlformats.org/package/2006/content-types}Override') if override.get('PartName') and '/slides/slide' in override.get('PartName') and not override.get('PartName').endswith('/slides/slide1.xml')]
        for override in overrides_to_remove:
            root.remove(override)
        tree.write(content_types_file, encoding='utf-8', xml_declaration=True)

def merge_single_slide_presentations(pptx_files, output_file):
    if not pptx_files:
        raise ValueError("Cannot merge an empty list of files.")
    with tempfile.TemporaryDirectory() as temp_dir:
        extracted_dirs = []
        for idx, pptx in enumerate(pptx_files):
            extract_path = os.path.join(temp_dir, f"pptx_{idx}")
            with zipfile.ZipFile(pptx, 'r') as zip_ref:
                zip_ref.extractall(extract_path)
            extracted_dirs.append(extract_path)
        base_dir = extracted_dirs[0]
        slides_dir = os.path.join(base_dir, "ppt", "slides")
        slides_rels_dir = os.path.join(base_dir, "ppt", "slides", "_rels")
        os.makedirs(slides_rels_dir, exist_ok=True)
        slide_num = 1
        for idx, extract_path in enumerate(extracted_dirs):
            if idx == 0:
                continue
            slide_num += 1
            src_slide = os.path.join(extract_path, "ppt", "slides", "slide1.xml")
            dst_slide = os.path.join(slides_dir, f"slide{slide_num}.xml")
            shutil.copy2(src_slide, dst_slide)
            src_rel = os.path.join(extract_path, "ppt", "slides", "_rels", "slide1.xml.rels")
            dst_rel = os.path.join(slides_rels_dir, f"slide{slide_num}.xml.rels")
            if os.path.exists(src_rel):
                shutil.copy2(src_rel, dst_rel)
        update_merged_presentation_structure(base_dir, slide_num)
        with zipfile.ZipFile(output_file, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
            for root, dirs, files in os.walk(base_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, base_dir)
                    zip_ref.write(file_path, arcname)

def update_merged_presentation_structure(base_dir, num_slides):
    pres_file = os.path.join(base_dir, "ppt", "presentation.xml")
    if os.path.exists(pres_file):
        tree = ET.parse(pres_file)
        root = tree.getroot()
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        slide_id_list = root.find('.//p:sldIdLst', ns)
        if slide_id_list is not None:
            slide_id_list.clear()
            for i in range(num_slides):
                slide_id = ET.SubElement(slide_id_list, '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
                slide_id.set('id', str(256 + i))
                slide_id.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', f'rId{i+2}')
        tree.write(pres_file, encoding='utf-8', xml_declaration=True)
    pres_rels_file = os.path.join(base_dir, "ppt", "_rels", "presentation.xml.rels")
    if os.path.exists(pres_rels_file):
        tree = ET.parse(pres_rels_file)
        root = tree.getroot()
        relationships_to_remove = [rel for rel in root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship') if rel.get('Type') == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide']
        for rel in relationships_to_remove:
            root.remove(rel)
        for i in range(num_slides):
            new_rel = ET.SubElement(root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
            new_rel.set('Id', f'rId{i+2}')
            new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
            new_rel.set('Target', f'slides/slide{i+1}.xml')
        tree.write(pres_rels_file, encoding='utf-8', xml_declaration=True)
    content_types_file = os.path.join(base_dir, "[Content_Types].xml")
    if os.path.exists(content_types_file):
        tree = ET.parse(content_types_file)
        root = tree.getroot()
        overrides_to_remove = [override for override in root.findall('.//{http://schemas.openxmlformats.org/package/2006/content-types}Override') if override.get('PartName') and '/slides/slide' in override.get('PartName')]
        for override in overrides_to_remove:
            root.remove(override)
        for i in range(num_slides):
            new_override = ET.SubElement(root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
            new_override.set('PartName', f'/ppt/slides/slide{i+1}.xml')
            new_override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
        tree.write(content_types_file, encoding='utf-8', xml_declaration=True)

# --- END: High-fidelity slide copy functions ---


# --- Main NiFi Script Logic ---
def main():
    flowFile = session.get()
    if not flowFile:
        return

    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # 1. Read incoming FlowFile into a temporary file
            with session.read(flowFile) as reader:
                incoming_bytes = reader.read()
            base_pptx_path = os.path.join(temp_dir, "base.pptx")
            with open(base_pptx_path, 'wb') as f:
                f.write(incoming_bytes)

            # 2. Check if the local file exists
            if not os.path.exists(LOCAL_PPTX_TO_APPEND):
                log.error(f"Local PPTX file not found at: {LOCAL_PPTX_TO_APPEND}")
                session.transfer(flowFile, REL_FAILURE)
                return

            # 3. Deconstruct both presentations into single-slide PPTX files
            all_single_slide_files = []

            # Process base presentation
            try:
                base_pres = Presentation(base_pptx_path)
                num_slides_base = len(base_pres.slides)
            except Exception as e:
                log.error(f"Could not open base presentation: {e}")
                session.transfer(flowFile, REL_FAILURE)
                return

            for i in range(num_slides_base):
                slide_path = os.path.join(temp_dir, f"base_slide_{i}.pptx")
                copy_single_slide(base_pptx_path, slide_path, i + 1)
                all_single_slide_files.append(slide_path)

            # Process append presentation
            try:
                append_pres = Presentation(LOCAL_PPTX_TO_APPEND)
                num_slides_append = len(append_pres.slides)
            except Exception as e:
                log.error(f"Could not open local presentation to append: {e}")
                session.transfer(flowFile, REL_FAILURE)
                return

            for i in range(num_slides_append):
                slide_path = os.path.join(temp_dir, f"append_slide_{i}.pptx")
                copy_single_slide(LOCAL_PPTX_TO_APPEND, slide_path, i + 1)
                all_single_slide_files.append(slide_path)

            # 4. Merge all the single-slide presentations into a final output file
            final_output_path = os.path.join(temp_dir, "merged.pptx")
            merge_single_slide_presentations(all_single_slide_files, final_output_path)

            # 5. Write the final merged file back to the FlowFile
            with open(final_output_path, 'rb') as f:
                final_bytes = f.read()
            session.write(flowFile, final_bytes)

            session.transfer(flowFile, REL_SUCCESS)

        except Exception as e:
            log.error(f"Failed to merge PPTX files with high fidelity: {e}")
            import traceback
            log.error(traceback.format_exc())
            session.transfer(flowFile, REL_FAILURE)

main()