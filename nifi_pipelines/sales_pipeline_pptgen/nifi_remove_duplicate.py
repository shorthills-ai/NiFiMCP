import sys
import io
import hashlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

class EfficientPPTDuplicateRemover:
    def __init__(self, pptx_bytes):
        self.presentation = Presentation(io.BytesIO(pptx_bytes))
        self.duplicate_slides = []

    def _get_slide_fingerprint(self, slide):
        """
        Create a lightweight fingerprint using key metrics
        Most efficient approach - Time complexity: O(n) where n is number of shapes
        """
        fingerprint = {
            'shape_count': len(slide.shapes),
            'text_char_count': 0,
            'image_count': 0,
            'shape_types': [],
            'positions': [],
            'text_hash': None
        }
        
        text_parts = []
        for shape in slide.shapes:
            # Count shape types
            fingerprint['shape_types'].append(int(shape.shape_type))
            
            # Store position as a simple tuple (rounded to nearest 1000 for fuzzy matching)
            pos = (round(shape.left/1000), round(shape.top/1000), 
                   round(shape.width/1000), round(shape.height/1000))
            fingerprint['positions'].append(pos)
            
            # Count images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                fingerprint['image_count'] += 1
            
            # Collect text efficiently
            if hasattr(shape, 'text') and shape.text.strip():
                text_parts.append(shape.text.strip())
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            text_parts.append(run.text.strip())
        
        # Create a quick hash of text content only if there's text
        if text_parts:
            combined_text = ''.join(text_parts)
            fingerprint['text_char_count'] = len(combined_text)
            fingerprint['text_hash'] = hashlib.md5(combined_text.encode('utf-8')).hexdigest()[:8]  # Short hash
        
        # Sort lists for consistent comparison
        fingerprint['shape_types'].sort()
        fingerprint['positions'].sort()
        
        return fingerprint
    
    def _slides_match_by_fingerprint(self, fp1, fp2):
        """
        Compare two fingerprints for similarity
        Allows for minor differences in positioning
        """
        # Must have same number of shapes
        if fp1['shape_count'] != fp2['shape_count']:
            return False
        
        # Must have same number of images
        if fp1['image_count'] != fp2['image_count']:
            return False
        
        # Must have same shape types
        if fp1['shape_types'] != fp2['shape_types']:
            return False
        
        # Text must match exactly if present
        if fp1['text_hash'] != fp2['text_hash']:
            return False
        
        # Positions should be similar (allowing for minor differences)
        if len(fp1['positions']) != len(fp2['positions']):
            return False
        
        # Sort positions and compare
        pos1_sorted = sorted(fp1['positions'])
        pos2_sorted = sorted(fp2['positions'])
        
        return pos1_sorted == pos2_sorted
    
    def analyze_and_remove(self):
        """Analyze all slides and identify duplicates using fingerprint method"""
        fingerprints = []
        for i, slide in enumerate(self.presentation.slides):
            current_fp = self._get_slide_fingerprint(slide)
            is_duplicate = False
            for j, existing_fp in enumerate(fingerprints):
                if self._slides_match_by_fingerprint(current_fp, existing_fp):
                    self.duplicate_slides.append(i)
                    is_duplicate = True
                    break
            if not is_duplicate:
                fingerprints.append(current_fp)
        # Remove duplicates in reverse order
        self.duplicate_slides.sort(reverse=True)
        for slide_idx in self.duplicate_slides:
            slide_element = self.presentation.slides._sldIdLst[slide_idx]
            self.presentation.slides._sldIdLst.remove(slide_element)

    def get_cleaned_pptx_bytes(self):
        out = io.BytesIO()
        self.presentation.save(out)
        return out.getvalue()

def main():
    pptx_bytes = sys.stdin.buffer.read()
    remover = EfficientPPTDuplicateRemover(pptx_bytes)
    remover.analyze_and_remove()
    cleaned_bytes = remover.get_cleaned_pptx_bytes()
    sys.stdout.buffer.write(cleaned_bytes)

if __name__ == "__main__":
    main()